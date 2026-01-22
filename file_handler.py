"""
file_handler.py - Extract text content from uploaded files
Supports: PDF, DOCX, PPTX, TXT
"""

import io
from typing import Optional, Dict
import PyPDF2
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_file = io.BytesIO(file_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_content = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text_content.append(page.extract_text())
        
        return "\n\n".join(text_content).strip()
    except Exception as e:
        raise ValueError(f"Error extracting PDF: {str(e)}")


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX file"""
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = Document(docx_file)
        
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        return "\n\n".join(text_content).strip()
    except Exception as e:
        raise ValueError(f"Error extracting DOCX: {str(e)}")


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX file"""
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        
        text_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content).strip()
    except Exception as e:
        raise ValueError(f"Error extracting PPTX: {str(e)}")


def extract_text_from_txt(file_bytes: bytes) -> str:
    """Extract text from TXT file"""
    try:
        return file_bytes.decode('utf-8').strip()
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            return file_bytes.decode('latin-1').strip()
        except Exception as e:
            raise ValueError(f"Error reading TXT file: {str(e)}")


def process_uploaded_file(uploaded_file) -> Dict[str, str]:
    """
    Process uploaded file and extract text content
    
    Args:
        uploaded_file: Streamlit UploadedFile object
        
    Returns:
        Dict with 'content', 'filename', 'file_type', and optional 'error'
    """
    result = {
        "filename": uploaded_file.name,
        "file_type": uploaded_file.type,
        "content": "",
        "error": None
    }
    
    try:
        file_bytes = uploaded_file.read()
        file_extension = uploaded_file.name.lower().split('.')[-1]
        
        # Check file size (max 10MB)
        max_size = 10 * 1024 * 1024  # 10MB in bytes
        if len(file_bytes) > max_size:
            result["error"] = "File size exceeds 10MB limit"
            return result
        
        # Extract text based on file type
        if file_extension == 'pdf':
            result["content"] = extract_text_from_pdf(file_bytes)
        elif file_extension == 'docx':
            result["content"] = extract_text_from_docx(file_bytes)
        elif file_extension == 'pptx':
            result["content"] = extract_text_from_pptx(file_bytes)
        elif file_extension in ['txt', 'text']:
            result["content"] = extract_text_from_txt(file_bytes)
        else:
            result["error"] = f"Unsupported file type: .{file_extension}"
            return result
        
        # Check if content was extracted
        if not result["content"]:
            result["error"] = "No text content found in file"
            return result
        
        # Limit content length (max 5000 chars for processing)
        if len(result["content"]) > 5000:
            result["content"] = result["content"][:5000] + "\n\n... (content truncated for processing)"
        
    except Exception as e:
        result["error"] = f"Error processing file: {str(e)}"
    
    return result


def create_file_based_prompt(file_content: str, file_type: str) -> str:
    """
    Create a contextual prompt based on file type
    
    Args:
        file_content: Extracted text from file
        file_type: Type of file (resume, presentation, etc.)
        
    Returns:
        Formatted prompt for post generation
    """
    
    # Detect content type
    content_lower = file_content.lower()
    
    if any(keyword in content_lower for keyword in ['resume', 'cv', 'experience', 'education', 'skills']):
        prompt_type = "resume"
    elif any(keyword in content_lower for keyword in ['slide', 'presentation', 'agenda']):
        prompt_type = "presentation"
    elif any(keyword in content_lower for keyword in ['report', 'analysis', 'findings', 'conclusion']):
        prompt_type = "report"
    else:
        prompt_type = "document"
    
    # Create contextual prompts
    prompts = {
        "resume": (
            f"Based on this resume/CV content, create a professional LinkedIn post highlighting:\n"
            f"- Key achievements and skills\n"
            f"- Career progression or milestones\n"
            f"- Professional value proposition\n\n"
            f"Content:\n{file_content[:2000]}"
        ),
        "presentation": (
            f"Based on this presentation, create a LinkedIn post that:\n"
            f"- Summarizes the key insights\n"
            f"- Highlights main takeaways\n"
            f"- Engages the audience with the core message\n\n"
            f"Content:\n{file_content[:2000]}"
        ),
        "report": (
            f"Based on this report, create a LinkedIn post that:\n"
            f"- Shares the most important findings\n"
            f"- Provides actionable insights\n"
            f"- Invites professional discussion\n\n"
            f"Content:\n{file_content[:2000]}"
        ),
        "document": (
            f"Based on this document, create an engaging LinkedIn post that:\n"
            f"- Captures the main ideas\n"
            f"- Adds professional context\n"
            f"- Encourages meaningful engagement\n\n"
            f"Content:\n{file_content[:2000]}"
        )
    }
    
    return prompts.get(prompt_type, prompts["document"])